import os
import json
from typing import List, Dict
from pptx import Presentation
import PyPDF2
import logging
from config.settings import settings

logger = logging.getLogger(__name__)

class ContentExtractor:
    """Extract text and metadata from PPT and PDF files"""
    
    def __init__(self):
        self.supported_formats = ['.pptx', '.ppt', '.pdf']
    
    def extract_content(self, file_path: str, run_id: str) -> Dict:
        """
        Extract content from presentation files
        Returns: Dict with slides, metadata, and extracted text
        """
        # Ensure run-specific directories are created
        settings.ensure_directories(run_id)
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        if file_ext in ['.pptx', '.ppt']:
            return self._extract_from_ppt(file_path, run_id)
        elif file_ext == '.pdf':
            return self._extract_from_pdf(file_path, run_id)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    def _extract_from_ppt(self, file_path: str, run_id: str) -> Dict:
        """Extract content from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            slides_content = []
            extracted_dir = settings.EXTRACTED_DIR
            
            if extracted_dir is None:
                raise ValueError("EXTRACTED_DIR is not set. Ensure directories are initialized.")
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_notes = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_notes = slide.notes_slide.notes_text_frame.text.strip()
                
                slide_data = {
                    'slide_number': i + 1,
                    'title': slide_text[0] if slide_text else f"Slide {i + 1}",
                    'content': '\n'.join(slide_text),
                    'notes': slide_notes,
                    'text_elements': slide_text
                }
                
                # Save slide content to JSON
                slide_filename = f"slide_{i+1:03d}.json"
                slide_path = os.path.join(extracted_dir, slide_filename)
                with open(slide_path, 'w', encoding='utf-8') as f:
                    json.dump(slide_data, f, indent=2)
                slide_data['content_path'] = slide_path
                
                slides_content.append(slide_data)
            
            result = {
                'file_type': 'ppt',
                'total_slides': len(slides_content),
                'slides': slides_content,
                'metadata': {
                    'title': getattr(prs.core_properties, 'title', '') or 'Untitled Presentation',
                    'author': getattr(prs.core_properties, 'author', '') or 'Unknown',
                    'subject': getattr(prs.core_properties, 'subject', '') or ''
                }
            }
            
            # Save full content data
            content_path = os.path.join(extracted_dir, "full_content.json")
            with open(content_path, 'w', encoding='utf-8') as f:
                json.dump(result, f, indent=2)
            result['content_path'] = content_path
            
            return result
            
        except Exception as e:
            logger.error(f"Error extracting PPT content: {str(e)}")
            raise
    
    def _extract_from_pdf(self, file_path: str, run_id: str) -> Dict:
        """Extract content from PDF files"""
        try:
            slides_content = []
            extracted_dir = settings.EXTRACTED_DIR
            
            if extracted_dir is None:
                raise ValueError("EXTRACTED_DIR is not set. Ensure directories are initialized.")
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for i, page in enumerate(pdf_reader.pages):
                    text = page.extract_text().strip()
                    
                    if text:
                        lines = text.split('\n')
                        title = lines[0] if lines else f"Page {i + 1}"
                        content = '\n'.join(lines[1:]) if len(lines) > 1 else text
                        
                        slide_data = {
                            'slide_number': i + 1,
                            'title': title[:100] + '...' if len(title) > 100 else title,
                            'content': content,
                            'notes': '',
                            'text_elements': lines
                        }
                        
                        # Save slide content to JSON
                        slide_filename = f"slide_{i+1:03d}.json"
                        slide_path = os.path.join(extracted_dir, slide_filename)
                        with open(slide_path, 'w', encoding='utf-8') as f:
                            json.dump(slide_data, f, indent=2)
                        slide_data['content_path'] = slide_path
                        
                        slides_content.append(slide_data)
                
                metadata = pdf_reader.metadata or {}
                
                result = {
                    'file_type': 'pdf',
                    'total_slides': len(slides_content),
                    'slides': slides_content,
                    'metadata': {
                        'title': metadata.get('/Title', 'Untitled Document'),
                        'author': metadata.get('/Author', 'Unknown'),
                        'subject': metadata.get('/Subject', '')
                    }
                }
                
                # Save full content data
                content_path = os.path.join(extracted_dir, "full_content.json")
                with open(content_path, 'w', encoding='utf-8') as f:
                    json.dump(result, f, indent=2)
                result['content_path'] = content_path
                
                return result
                
        except Exception as e:
            logger.error(f"Error extracting PDF content: {str(e)}")
            raise